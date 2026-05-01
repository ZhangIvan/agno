import asyncio
import os
from pathlib import Path
from typing import IO, Any, List, Optional, Union
from uuid import uuid4

from agno.knowledge.chunking.document import DocumentChunking
from agno.knowledge.chunking.strategy import ChunkingStrategy, ChunkingStrategyType
from agno.knowledge.document.base import Document
from agno.knowledge.reader.base import Reader
from agno.knowledge.types import ContentType
from agno.utils.log import log_debug, log_error, log_warning

try:
    from pptx import Presentation  # type: ignore
except ImportError:
    raise ImportError("The `python-pptx` package is not installed. Please install it via `pip install python-pptx`.")


class PPTXReader(Reader):
    """Reader for PPTX files"""

    def __init__(
        self,
        chunking_strategy: Optional[ChunkingStrategy] = None,
        capture_pages: bool = True,
        pages_cache_dir: Optional[str] = None,
        image_dpi: int = 100,
        **kwargs,
    ):
        if chunking_strategy is None:
            chunk_size = kwargs.get("chunk_size", 5000)
            chunking_strategy = DocumentChunking(chunk_size=chunk_size)
        self.capture_pages = capture_pages
        self.pages_cache_dir = pages_cache_dir or os.getenv("AGNO_PAGE_CACHE_DIR")
        self.image_dpi = image_dpi
        super().__init__(chunking_strategy=chunking_strategy, **kwargs)

    @classmethod
    def get_supported_chunking_strategies(cls) -> List[ChunkingStrategyType]:
        """Get the list of supported chunking strategies for PPTX readers."""
        return [
            ChunkingStrategyType.DOCUMENT_CHUNKER,
            ChunkingStrategyType.CODE_CHUNKER,
            ChunkingStrategyType.FIXED_SIZE_CHUNKER,
            ChunkingStrategyType.SEMANTIC_CHUNKER,
            ChunkingStrategyType.AGENTIC_CHUNKER,
            ChunkingStrategyType.RECURSIVE_CHUNKER,
        ]

    @classmethod
    def get_supported_content_types(cls) -> List[ContentType]:
        return [ContentType.PPTX]

    def read(self, file: Union[Path, IO[Any]], name: Optional[str] = None) -> List[Document]:
        """Read a pptx file and return a list of documents (one per slide)."""
        _tmp_capture_path: Optional[str] = None
        try:
            file_path: Optional[str] = None
            if isinstance(file, Path):
                if not file.exists():
                    raise FileNotFoundError(f"Could not find file: {file}")
                log_debug(f"Reading: {file}")
                presentation = Presentation(str(file))
                doc_name = name or file.stem
                file_path = str(file)
            else:
                log_debug(f"Reading uploaded file: {getattr(file, 'name', 'BytesIO')}")
                presentation = Presentation(file)
                doc_name = name or getattr(file, "name", "pptx_file").split(".")[0]
                if self.capture_pages:
                    import os
                    import shutil
                    import tempfile

                    if hasattr(file, "seek"):
                        file.seek(0)
                    _tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
                    shutil.copyfileobj(file, _tmp)
                    _tmp.close()
                    if hasattr(file, "seek"):
                        file.seek(0)
                    file_path = _tmp.name
                    _tmp_capture_path = _tmp.name

            total_slides = len(presentation.slides)

            # Capture slide images if requested (file path required)
            page_images: Optional[dict] = None
            if self.capture_pages and file_path:
                try:
                    from agno.knowledge.reader.page_capture import capture_pptx_slides, get_page_cache_dir

                    cache_dir = get_page_cache_dir(self.pages_cache_dir, doc_name)
                    page_images = capture_pptx_slides(file_path, cache_dir, dpi=self.image_dpi)
                except Exception as e:
                    log_warning(f"Failed to capture PPTX slide images: {e}")

            # When page images were captured, use image-only mode:
            # slides are visual artefacts — text extraction is skipped and only
            # page_image documents are produced for multimodal embedding.
            # Fallback to text extraction when capture was unavailable (e.g. no LibreOffice).
            if page_images:
                result: List[Document] = []
                for slide_num, image_path in page_images.items():
                    img_meta: dict = {
                        "doc_type": "page_image",
                        "page_number": slide_num,
                        "total_pages": total_slides,
                        "page_image_path": image_path,
                    }
                    if cache_dir:  # type: ignore[possibly-undefined]
                        img_meta["pages_cache_dir"] = cache_dir
                    result.append(
                        Document(
                            name=doc_name,
                            id=f"{doc_name}_img_{slide_num}",
                            content="",
                            meta_data=img_meta,
                            content_id=f"{doc_name}_page_{slide_num}",
                        )
                    )
            else:
                # Text-only fallback (LibreOffice unavailable)
                documents: List[Document] = []
                for slide_number, slide in enumerate(presentation.slides, 1):
                    text_content = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            text_content.append(shape.text.strip())

                    slide_text = f"Slide {slide_number}:\n"
                    slide_text += "\n".join(text_content) if text_content else "(No text content)"

                    documents.append(
                        Document(
                            name=doc_name,
                            id=str(uuid4()),
                            meta_data={
                                "page_number": slide_number,
                                "total_pages": total_slides,
                                "doc_type": "text_chunk",
                            },
                            content=slide_text,
                        )
                    )

                if self.chunk:
                    chunked_documents = []
                    for document in documents:
                        chunked_documents.extend(self.chunk_document(document))
                    result = chunked_documents
                else:
                    result = documents

            return result

        except Exception as e:
            log_error(f"Error reading file: {e}")
            raise ValueError(f"Error reading file: {e}")
        finally:
            if _tmp_capture_path:
                try:
                    import os

                    os.unlink(_tmp_capture_path)
                except OSError:
                    pass

    async def async_read(self, file: Union[Path, IO[Any]], name: Optional[str] = None) -> List[Document]:
        """Asynchronously read a pptx file and return a list of documents"""
        try:
            return await asyncio.to_thread(self.read, file, name)
        except Exception as e:
            log_error(f"Error reading file asynchronously: {e}")
            raise ValueError(f"Error reading file asynchronously: {e}")
